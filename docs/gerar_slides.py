# -*- coding: utf-8 -*-
"""Gera os slides da apresentacao (Fase 4) em PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

AZUL = RGBColor(0x1B, 0x3A, 0x5C)
VERDE = RGBColor(0x2E, 0x7D, 0x32)
CINZA = RGBColor(0x4A, 0x4A, 0x4A)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
LARGURA = prs.slide_width
ALTURA = prs.slide_height
BRANCO_LAYOUT = prs.slide_layouts[6]  # layout em branco


def nova_slide():
    return prs.slides.add_slide(BRANCO_LAYOUT)


def fundo(slide, cor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = cor


def caixa_texto(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def titulo_secao(slide, texto, cor_fundo=AZUL, cor_texto=BRANCO):
    fundo(slide, cor_fundo)
    tf = caixa_texto(slide, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.8))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = cor_texto


def slide_conteudo(titulo, bullets, subtitulo=None, tamanho_titulo=32, tamanho_bullet=20):
    slide = nova_slide()
    fundo(slide, BRANCO)

    # barra superior
    barra = slide.shapes.add_shape(1, 0, 0, LARGURA, Inches(0.15))
    barra.fill.solid()
    barra.fill.fore_color.rgb = VERDE
    barra.line.fill.background()

    tf_titulo = caixa_texto(slide, Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    p = tf_titulo.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(tamanho_titulo)
    p.font.bold = True
    p.font.color.rgb = AZUL

    if subtitulo:
        tf_sub = caixa_texto(slide, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.6))
        p = tf_sub.paragraphs[0]
        p.text = subtitulo
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = CINZA
        top_bullets = Inches(1.9)
    else:
        top_bullets = Inches(1.6)

    tf = caixa_texto(slide, Inches(0.8), top_bullets, Inches(11.7), Inches(5.2))
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        nivel = 0
        texto_item = item
        if isinstance(item, tuple):
            texto_item, nivel = item
        p.text = ("•  " if nivel == 0 else "     -  ") + texto_item
        p.font.size = Pt(tamanho_bullet - (4 if nivel else 0))
        p.font.color.rgb = CINZA if nivel else RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(14 if nivel == 0 else 6)
    return slide


# ---------------------------------------------------------------------
# 1. Capa
# ---------------------------------------------------------------------
slide = nova_slide()
fundo(slide, AZUL)
tf = caixa_texto(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.0))
p = tf.paragraphs[0]
p.text = "AgroSele"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = BRANCO
p2 = tf.add_paragraph()
p2.text = "Seleção Automática de Respostas Técnicas para Produtores Rurais\nde Pecuária Leiteira em Português"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0xC8, 0xD8, 0xE8)
p2.space_before = Pt(16)

tf2 = caixa_texto(slide, Inches(0.9), Inches(6.2), Inches(11.5), Inches(1.0))
p = tf2.paragraphs[0]
p.text = "Frederico Botelho Martins  |  Processamento de Linguagem Natural  |  Prof. Dr. Yuri Malheiros  |  2026.1"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0xA0, 0xB8, 0xD0)

# ---------------------------------------------------------------------
# 2. Motivacao
# ---------------------------------------------------------------------
slide_conteudo(
    "Motivação",
    [
        "Serviços de extensão rural (ex.: Embrapa Gado de Leite) recebem milhares de perguntas técnicas de produtores ao longo dos anos",
        "Muitas perguntas são semanticamente equivalentes a outras já respondidas por especialistas",
        "Hoje, cada resposta exige trabalho manual e individualizado, um modelo que não escala",
        "Domínio desafiador: perguntas informais e heterogêneas, vocabulário técnico compartilhado entre respostas de assuntos distintos",
        "Motivação pessoal: o tema conecta-se a um software de gestão para fazendas leiteiras já em produção, desenvolvido pelo autor",
    ],
)

# ---------------------------------------------------------------------
# 3. Pergunta de pesquisa
# ---------------------------------------------------------------------
slide = nova_slide()
fundo(slide, RGBColor(0xF4, 0xF7, 0xFA))
barra = slide.shapes.add_shape(1, 0, 0, LARGURA, Inches(0.15))
barra.fill.solid()
barra.fill.fore_color.rgb = VERDE
barra.line.fill.background()
tf_titulo = caixa_texto(slide, Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.9))
p = tf_titulo.paragraphs[0]
p.text = "Pergunta de Pesquisa"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = AZUL

caixa = slide.shapes.add_shape(1, Inches(1.2), Inches(2.3), Inches(11.0), Inches(3.2))
caixa.fill.solid()
caixa.fill.fore_color.rgb = BRANCO
caixa.line.color.rgb = VERDE
caixa.line.width = Pt(2)
tf = caixa.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.5)
tf.margin_right = Inches(0.5)
p = tf.paragraphs[0]
p.text = ("Até que ponto diferentes estratégias de combinação entre sinal lexical clássico "
          "e representações neurais pré-treinadas, congeladas, fundidas ou ajustadas por "
          "fine-tuning, afetam o desempenho e o custo computacional da seleção de resposta "
          "em um domínio técnico de baixo recurso como a pecuária leiteira em português?")
p.font.size = Pt(24)
p.font.italic = True
p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
p.alignment = PP_ALIGN.CENTER

# ---------------------------------------------------------------------
# 4. Dataset
# ---------------------------------------------------------------------
slide_conteudo(
    "Dados: MilkQA",
    [
        "Dataset real (Criscuolo et al., 2018): atendimento ao produtor da Embrapa Gado de Leite, 2003-2012",
        "2.657 pares pergunta-resposta, escritos por produtores reais e respondidos por especialistas",
        "Tarefa: dado 1 pergunta + 50 respostas candidatas (1 correta, 49 distratoras), apontar a correta",
        "Splits oficiais: treino (2.307) | validação (50) | teste (300)",
        "Perguntas longas e informais (média de 198 palavras, até 682)",
    ],
    subtitulo="Seleção de resposta (answer selection) — não é classificação em classes fixas, é ranqueamento",
)

# ---------------------------------------------------------------------
# 5. Abordagem -- visao geral das 6 variantes
# ---------------------------------------------------------------------
slide_conteudo(
    "Abordagem: 6 configurações comparadas",
    [
        "1. TF-IDF clássico — tokenização + cosseno, implementado do zero, sem BERT",
        "2. Bi-encoder — BERTimbau congelado + MLP sobre [q, a, |q−a|, q*a]",
        "3. Cross-Encoder — atenção cruzada entre pergunta e candidata, congelado",
        "4. Híbrido — BERTimbau congelado + escore BM25 fundido no vetor de entrada",
        "5. Pipeline multi-estágio — BM25 filtra top-K, Cross-Encoder reranqueia",
        "6. Fine-tuning parcial — última camada do BERTimbau descongelada e ajustada",
    ],
    subtitulo="Do método mais simples ao mais sofisticado, sobre o mesmo protocolo de avaliação",
)

# ---------------------------------------------------------------------
# 6. Metricas
# ---------------------------------------------------------------------
slide_conteudo(
    "Métricas de Avaliação",
    [
        "Accuracy@1 — a resposta correta ficou em 1º lugar no ranking das 50 candidatas?",
        "MRR (Mean Reciprocal Rank) — média do inverso da posição da resposta correta",
        "F1-macro não se aplica: a tarefa é de ranqueamento, não classificação em classes fixas",
        "Baseline aleatório de referência: Accuracy@1 ≈ 0,020 | MRR ≈ 0,090",
    ],
)

# ---------------------------------------------------------------------
# 7. Resultados -- grafico de barras nativo
# ---------------------------------------------------------------------
slide = nova_slide()
fundo(slide, BRANCO)
barra = slide.shapes.add_shape(1, 0, 0, LARGURA, Inches(0.15))
barra.fill.solid()
barra.fill.fore_color.rgb = VERDE
barra.line.fill.background()
tf_titulo = caixa_texto(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.8))
p = tf_titulo.paragraphs[0]
p.text = "Resultados (Accuracy@1, conjunto de teste — 300 perguntas)"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = AZUL

chart_data = CategoryChartData()
chart_data.categories = [
    "Aleatório", "Cosseno\npuro", "TF-IDF", "Bi-encoder\ncongelado",
    "Cross-\nEncoder", "Pipeline\n(K=10)", "Híbrido\nBM25", "Fine-tuning\nparcial",
]
chart_data.add_series("Accuracy@1", (0.020, 0.277, 0.503, 0.570, 0.617, 0.590, 0.663, 0.690))

x, y, cx, cy = Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6)
grafico = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
grafico.has_legend = False
plot = grafico.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.000'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(12)
serie = plot.series[0]
serie.format.fill.solid()
serie.format.fill.fore_color.rgb = VERDE
cat_axis = grafico.category_axis
cat_axis.tick_labels.font.size = Pt(11)
val_axis = grafico.value_axis
val_axis.tick_labels.font.size = Pt(11)
val_axis.maximum_scale = 0.8

# ---------------------------------------------------------------------
# 8. Achado principal
# ---------------------------------------------------------------------
slide_conteudo(
    "Achado Principal",
    [
        "TF-IDF (sem BERT nenhum) quase dobra o BERTimbau cru: 0,503 vs 0,277",
        ("→ vocabulário técnico compartilhado favorece correspondência lexical exata", 1),
        "Fusão híbrida (BERT congelado + BM25) chega a 0,663 sem tocar nos pesos do BERT",
        ("→ recupera boa parte do ganho do fine-tuning, a 2 min vs 11h de treino", 1),
        "Fine-tuning parcial (0,690) é o melhor resultado isolado — mas só depois que uma primeira tentativa com dado insuficiente (300 de 2.307 perguntas) sugeriu o contrário",
        ("→ limitação computacional pode mascarar o real potencial de uma técnica", 1),
    ],
)

# ---------------------------------------------------------------------
# 9. Exemplo concreto (entrada/saida)
# ---------------------------------------------------------------------
slide = nova_slide()
fundo(slide, RGBColor(0xF4, 0xF7, 0xFA))
barra = slide.shapes.add_shape(1, 0, 0, LARGURA, Inches(0.15))
barra.fill.solid()
barra.fill.fore_color.rgb = VERDE
barra.line.fill.background()
tf_titulo = caixa_texto(slide, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.7))
p = tf_titulo.paragraphs[0]
p.text = "Exemplo Concreto: Entrada e Saída do Sistema"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = AZUL

caixa_p = slide.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(12.1), Inches(1.3))
caixa_p.fill.solid()
caixa_p.fill.fore_color.rgb = RGBColor(0xE8, 0xEF, 0xF5)
caixa_p.line.color.rgb = AZUL
tf = caixa_p.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "ENTRADA (pergunta do produtor):"
p.font.bold = True
p.font.size = Pt(14)
p.font.color.rgb = AZUL
p2 = tf.add_paragraph()
p2.text = ("\"Meu tio Paulo mora em uma chácara... o leite azeda muito rápido... "
           "o que devemos fazer?\" + 50 respostas candidatas do corpus")
p2.font.size = Pt(15)
p2.font.italic = True

caixa_s = slide.shapes.add_shape(1, Inches(0.6), Inches(2.8), Inches(12.1), Inches(1.6))
caixa_s.fill.solid()
caixa_s.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xE6)
caixa_s.line.color.rgb = VERDE
tf = caixa_s.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "SAÍDA (resposta selecionada pelo sistema, rank 1º de 50):"
p.font.bold = True
p.font.size = Pt(14)
p.font.color.rgb = VERDE
p2 = tf.add_paragraph()
p2.text = ("Resposta técnica da Embrapa sobre conservação do leite: cuidados de higiene na "
           "ordenha, resfriamento imediato e tempo máximo antes do processamento — "
           "correspondência correta confirmada pelo gabarito do MilkQA.")
p2.font.size = Pt(15)

caixa_e = slide.shapes.add_shape(1, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.2))
caixa_e.fill.solid()
caixa_e.fill.fore_color.rgb = RGBColor(0xFB, 0xEA, 0xE5)
caixa_e.line.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
tf = caixa_e.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "Caso de ERRO (análise qualitativa, fine-tuning parcial):"
p.font.bold = True
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
p2 = tf.add_paragraph()
p2.text = ('Pergunta sobre frequência/horário de alimentação das vacas → sistema escolheu '
           '(score 0,388) uma resposta com fórmula de ração de baixo custo, em vez da resposta '
           'correta sobre produção a pasto (score 0,000, posição 38/50).')
p2.font.size = Pt(14)
p3 = tf.add_paragraph()
p3.text = ('Causa provável: a resposta errada tem alta densidade de vocabulário técnico '
           '(insumos, proporções numéricas) que funciona como "atrator" lexical, mesmo '
           'divergindo do subtema exato da pergunta.')
p3.font.size = Pt(14)
p3.font.italic = True

# ---------------------------------------------------------------------
# 10. Limitacoes
# ---------------------------------------------------------------------
slide_conteudo(
    "Limitações",
    [
        "Truncamento de texto em 256 tokens (128 nas variantes com fine-tuning/Cross-Encoder) descarta informação de cauda em respostas longas",
        "Amostragem de negativos majoritariamente aleatória (hard negative mining testado, sem ganho estatístico)",
        "Backbone congelado nas variantes híbrida e Cross-Encoder — não explora combinar arquitetura + fine-tuning simultaneamente",
        "Sem comparação direta com uma API comercial de LLM (decisão metodológica: custo, latência, privacidade, determinismo)",
        "Modelo ainda erra sistematicamente diante de candidatas lexicalmente densas, mesmo quando fora do subtema exato",
    ],
)

# ---------------------------------------------------------------------
# 11. Conclusao / obrigado
# ---------------------------------------------------------------------
slide = nova_slide()
fundo(slide, AZUL)
tf = caixa_texto(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.2))
p = tf.paragraphs[0]
p.text = "Conclusão"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BRANCO

tf2 = caixa_texto(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(3.5))
itens = [
    "Sinal lexical explícito e volume de dado de treino pesam tanto quanto a arquitetura escolhida",
    "Fine-tuning parcial é o melhor resultado isolado (0,690), mas o híbrido BM25 chega perto (0,663) a uma fração do custo",
    "Pipeline multi-estágio recupera qualidade do Cross-Encoder com até 80% menos custo computacional",
    "Código, notebooks executados e artigo completo: github.com/fredmartins12/AgroSele-MilkQA",
]
for i, item in enumerate(itens):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = "•  " + item
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xD8, 0xE4, 0xF0)
    p.space_after = Pt(18)

tf3 = caixa_texto(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.6))
p = tf3.paragraphs[0]
p.text = "Obrigado! Perguntas?"
p.font.size = Pt(22)
p.font.italic = True
p.font.color.rgb = BRANCO

prs.save("Apresentacao_AgroSele.pptx")
print(f"slides gerados: {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
print("salvo em Apresentacao_AgroSele.pptx")
